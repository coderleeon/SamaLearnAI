"""PowerPoint presentation parser using python-pptx.

Extracts text from slides, keeping slide number metadata for citation formatting.
"""

import io
from pptx import Presentation


def parse_pptx(file_bytes: bytes, filename: str) -> list[dict]:
    """Extract text from a PowerPoint PPTX file slide by slide.

    Args:
        file_bytes: Raw bytes of the PPTX file.
        filename: Original filename for metadata.

    Returns:
        List of dicts with keys: content, metadata.
        metadata includes: source_type, filename, slide_number, total_slides.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    total_slides = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        text_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_parts.append(run.text)

        slide_text = " ".join(text_parts).strip()
        if not slide_text:
            continue

        slides.append({
            "content": slide_text,
            "metadata": {
                "source_type": "pptx",
                "filename": filename,
                "slide_number": idx + 1,  # 1-indexed for display
                "total_slides": total_slides,
            },
        })

    return slides
